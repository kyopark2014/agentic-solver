import json
import boto3
import os
import time
import datetime
import PyPDF2
import csv
import traceback
import re
import base64
import requests
import operator

from io import BytesIO
from urllib import parse
from botocore.config import Config
from PIL import Image
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.memory import ConversationBufferWindowMemory

from langchain_community.docstore.document import Document
from langchain_community.vectorstores.faiss import FAISS
from langchain_community.embeddings import BedrockEmbeddings
from multiprocessing import Process, Pipe

from langchain_core.prompts import MessagesPlaceholder, ChatPromptTemplate
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langchain_aws import ChatBedrock
from langchain_core.prompts import PromptTemplate

from langchain.agents import tool
from bs4 import BeautifulSoup
from pytz import timezone
from opensearchpy import OpenSearch, RequestsHttpConnection, AWSV4SignerAuth

from langchain_aws import AmazonKnowledgeBasesRetriever
from pydantic.v1 import BaseModel, Field
from tavily import TavilyClient  
from langgraph.graph import START, END, StateGraph
from typing import Annotated, List, Tuple, TypedDict, Literal, Sequence, Union
from langgraph.graph.message import add_messages
from langgraph.prebuilt import ToolNode

s3 = boto3.client('s3')
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
s3_arn = os.environ.get('s3_arn')
callLogTableName = os.environ.get('callLogTableName')

enableReference = os.environ.get('enableReference', 'false')
debugMessageMode = os.environ.get('debugMessageMode', 'false')
opensearch_url = os.environ.get('opensearch_url')
path = os.environ.get('path')
doc_prefix = s3_prefix+'/'
speech_prefix = 'speech/'
LLM_for_chat = json.loads(os.environ.get('LLM_for_chat'))
LLM_for_multimodal= json.loads(os.environ.get('LLM_for_multimodal'))
LLM_embedding = json.loads(os.environ.get('LLM_embedding'))
priorty_search_embedding = json.loads(os.environ.get('priorty_search_embedding'))
knowledge_base_name = os.environ.get('knowledge_base_name')
knowledge_base_role = os.environ.get('knowledge_base_role')
embeddingModelArn = os.environ.get('embeddingModelArn')
parsingModelArn = os.environ.get('parsingModelArn')
collectionArn = os.environ.get('collectionArn')
vectorIndexName = os.environ.get('vectorIndexName')

selected_chat = 0
selected_multimodal = 0
selected_embedding = 0
selected_ps_embedding = 0
length_of_models = 1

useParallelRAG = os.environ.get('useParallelRAG', 'true')
roleArn = os.environ.get('roleArn')
numberOfDocs = int(os.environ.get('numberOfRelevantDocs', '8'))
MSG_LENGTH = 100
MSG_HISTORY_LENGTH = 20
speech_generation = True
history_length = 0
token_counter_history = 0
grade_state = "LLM" # LLM, PRIORITY_SEARCH, OTHERS
 
minDocSimilarity = 350
projectName = os.environ.get('projectName')
maxOutputTokens = 4096
data_source_id = ""
multi_region_mode = "enable"  # CSAT

multi_region_models = [   # claude sonnet 3.0
    {   
        "bedrock_region": "us-west-2", # Oregon
        "model_type": "claude3",
        "model_id": "anthropic.claude-3-sonnet-20240229-v1:0"
    },
    {
        "bedrock_region": "us-east-1", # N.Virginia
        "model_type": "claude3",
        "model_id": "anthropic.claude-3-sonnet-20240229-v1:0"
    },
    {
        "bedrock_region": "ca-central-1", # Canada
        "model_type": "claude3",
        "model_id": "anthropic.claude-3-sonnet-20240229-v1:0"
    },
    {
        "bedrock_region": "eu-west-2", # London
        "model_type": "claude3",
        "model_id": "anthropic.claude-3-sonnet-20240229-v1:0"
    },
    {
        "bedrock_region": "sa-east-1", # Sao Paulo
        "model_type": "claude3",
        "model_id": "anthropic.claude-3-sonnet-20240229-v1:0"
    }
]
multi_region = 'disable'

reference_docs = []

secretsmanager = boto3.client('secretsmanager')

tavily_api_key = ""
weather_api_key = ""
def load_secrets():
    global tavily_api_key, weather_api_key
    secretsmanager = boto3.client('secretsmanager')
    
    # api key to use LangSmith
    langsmith_api_key = ""
    try:
        get_langsmith_api_secret = secretsmanager.get_secret_value(
            SecretId=f"langsmithapikey-{projectName}"
        )
        # print('get_langsmith_api_secret: ', get_langsmith_api_secret)
        
        secret = json.loads(get_langsmith_api_secret['SecretString'])
        #print('secret: ', secret)
        langsmith_api_key = secret['langsmith_api_key']
        langchain_project = secret['langchain_project']
    except Exception as e:
        raise e

    if langsmith_api_key:
        os.environ["LANGCHAIN_API_KEY"] = langsmith_api_key
        os.environ["LANGCHAIN_TRACING_V2"] = "true"
        os.environ["LANGCHAIN_PROJECT"] = langchain_project
        
    # api key to use Tavily Search    
    try:
        get_tavily_api_secret = secretsmanager.get_secret_value(
            SecretId=f"tavilyapikey-{projectName}"
        )
        #print('get_tavily_api_secret: ', get_tavily_api_secret)
        
        secret = json.loads(get_tavily_api_secret['SecretString'])
        # print('secret: ', secret)
        if secret['tavily_api_key']:
            tavily_api_key = json.loads(secret['tavily_api_key'])
        # print('tavily_api_key: ', tavily_api_key)
    except Exception as e: 
        raise e
    
    try:
        get_weather_api_secret = secretsmanager.get_secret_value(
            SecretId=f"openweathermap-{projectName}"
        )
        #print('get_weather_api_secret: ', get_weather_api_secret)
        
        secret = json.loads(get_weather_api_secret['SecretString'])
        #print('secret: ', secret)
        weather_api_key = secret['weather_api_key']

    except Exception as e:
        raise e
load_secrets()

def check_tavily_secret(tavily_api_key):
    query = 'what is LangGraph'
    valid_keys = []
    for key in tavily_api_key:
        try:
            tavily_client = TavilyClient(api_key=key)
            response = tavily_client.search(query, max_results=1)
            # print('tavily response: ', response)
            
            if 'results' in response and len(response['results']):
                valid_keys.append(key)
        except Exception as e:
            print('Exception: ', e)
    # print('valid_keys: ', valid_keys)
    
    return valid_keys

tavily_api_key = check_tavily_secret(tavily_api_key)
print('The number of valid tavily api keys: ', len(tavily_api_key))

selected_tavily = -1
if len(tavily_api_key):
    os.environ["TAVILY_API_KEY"] = tavily_api_key[0]
    selected_tavily = 0

def tavily_search(query, k):
    global selected_tavily
    docs = []
        
    if selected_tavily != -1:
        selected_tavily = selected_tavily + 1
        if selected_tavily == len(tavily_api_key):
            selected_tavily = 0

        try:
            tavily_client = TavilyClient(api_key=tavily_api_key[selected_tavily])
            response = tavily_client.search(query, max_results=k)
            
            # print('tavily response: ', response)
            
            for i, r in enumerate(response["results"]):
                contnet = r.get("content")
                print(f"{i}: {contnet}")
                
                name = r.get("title")
                if name is None:
                    name = 'WWW'
                
                url = ""
                if "url" in r:
                    url = r.get("url")
            
                docs.append(
                    Document(
                        page_content=contnet,
                        metadata={
                            'name': name,
                            'url': url,
                            'from': 'tavily'
                        },
                    )
                )   
        except Exception as e:
            print('Exception: ', e)
    return docs

#result = tavily_search('what is LangChain', 2)
#print('search result: ', result)

def get_multi_region_chat(models, selected):
    profile = models[selected]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    maxOutputTokens = 4096
    print(f'selected_chat: {selected}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )        
    return chat

weather_api_key = ""
def load_secret():
    global weather_api_key
    # api key to get weather information in agent    
    secretsmanager = boto3.client('secretsmanager')
    try:
        get_weather_api_secret = secretsmanager.get_secret_value(
            SecretId=f"openweathermap-{projectName}"
        )
        #print('get_weather_api_secret: ', get_weather_api_secret)
        secret = json.loads(get_weather_api_secret['SecretString'])
        #print('secret: ', secret)
        weather_api_key = secret['weather_api_key']

    except Exception as e:
        raise e
    
    # api key to use LangSmith
    langsmith_api_key = ""
    try:
        get_langsmith_api_secret = secretsmanager.get_secret_value(
            SecretId=f"langsmithapikey-{projectName}"
        )
        #print('get_langsmith_api_secret: ', get_langsmith_api_secret)
        secret = json.loads(get_langsmith_api_secret['SecretString'])
        #print('secret: ', secret)
        langsmith_api_key = secret['langsmith_api_key']
        langchain_project = secret['langchain_project']
    except Exception as e:
        raise e

    if langsmith_api_key:
        os.environ["LANGCHAIN_API_KEY"] = langsmith_api_key
        os.environ["LANGCHAIN_TRACING_V2"] = "true"
        os.environ["LANGCHAIN_PROJECT"] = langchain_project
        
    # api key to use Tavily Search
    tavily_api_key = ""
    try:
        get_tavily_api_secret = secretsmanager.get_secret_value(
            SecretId=f"tavilyapikey-{projectName}"
        )
        #print('get_tavily_api_secret: ', get_tavily_api_secret)
        secret = json.loads(get_tavily_api_secret['SecretString'])
        #print('secret: ', secret)
        tavily_api_key = secret['tavily_api_key']
    except Exception as e: 
        raise e

    if tavily_api_key:
        os.environ["TAVILY_API_KEY"] = tavily_api_key

load_secret()
    
# websocket
connection_url = os.environ.get('connection_url')
client = boto3.client('apigatewaymanagementapi', endpoint_url=connection_url)
print('connection_url: ', connection_url)

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"

map_chain = dict() 

def get_chat():
    global selected_chat
    
    if multi_region == 'enable':
        profile = multi_region_models[selected_chat]
    else:
        profile = LLM_for_chat[selected_chat]
    
    #print('length_of_models: ', length_of_models)
    #print('profile: ', json.dumps(profile))
        
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    maxOutputTokens = 4096
    print(f'LLM: {selected_chat}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    selected_chat = selected_chat + 1
    if selected_chat == length_of_models:
        selected_chat = 0
    
    return chat

def get_multimodal():
    global selected_multimodal
    
    profile = LLM_for_multimodal[selected_multimodal]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_multimodal}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    multimodal = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    selected_multimodal = selected_multimodal + 1
    if selected_multimodal == len(LLM_for_multimodal):
        selected_multimodal = 0
    
    return multimodal

def get_embedding():
    global selected_embedding
    profile = LLM_embedding[selected_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'Embedding: {selected_embedding}, bedrock_region: {bedrock_region}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_embedding = selected_embedding + 1
    if selected_embedding == len(LLM_embedding):
        selected_embedding = 0
    
    return bedrock_embedding

def get_ps_embedding():
    global selected_ps_embedding
    profile = priorty_search_embedding[selected_ps_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'selected_ps_embedding: {selected_ps_embedding}, bedrock_region: {bedrock_region}, model_id: {model_id}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region, 
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_ps_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_ps_embedding = selected_ps_embedding + 1
    if selected_ps_embedding == len(priorty_search_embedding):
        selected_ps_embedding = 0
    
    return bedrock_ps_embedding

def sendMessage(id, body):
    try:
        client.post_to_connection(
            ConnectionId=id, 
            Data=json.dumps(body)
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('err_msg: ', err_msg)
        # raise Exception ("Not able to send a message")

def sendResultMessage(connectionId, requestId, msg):    
    result = {
        'request_id': requestId,
        'msg': msg,
        'status': 'completed'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, result)

def sendDebugMessage(connectionId, requestId, msg):
    debugMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'debug'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, debugMsg)

def sendErrorMessage(connectionId, requestId, msg):
    errorMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'error'
    }
    print('error: ', json.dumps(errorMsg))
    sendMessage(connectionId, errorMsg)

def isKorean(text):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(text))
    # print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        print('Korean: ', word_kor)
        return True
    else:
        print('Not Korean: ', word_kor)
        return False

def general_conversation(connectionId, requestId, chat, query):
    global time_for_inference, history_length, token_counter_history    
    time_for_inference = history_length = token_counter_history = 0
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
    
    if isKorean(query)==True :
        system = (
            "다음의 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다."
        )
    else: 
        system = (
            "Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor."
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    # print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    # print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        isTyping(connectionId, requestId, "")  
        stream = chain.invoke(
            {
                "history": history,
                "input": query,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
                            
        msg = stream.content
        print('msg: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':  
        chat_history = ""
        for dialogue_turn in history:
            #print('type: ', dialogue_turn.type)
            #print('content: ', dialogue_turn.content)
            
            dialog = f"{dialogue_turn.type}: {dialogue_turn.content}\n"            
            chat_history = chat_history + dialog
                
        history_length = len(chat_history)
        print('chat_history length: ', history_length)
        
        token_counter_history = 0
        if chat_history:
            token_counter_history = chat.get_num_tokens(chat_history)
            print('token_size of history: ', token_counter_history)
        
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
        
    return msg
    
def get_summary(chat, docs):    
    text = ""
    for doc in docs:
        text = text + doc
    
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장을 요약해서 500자 이내로 설명하세오."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Write a concise summary within 500 characters."
        )
    
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        summary = result.content
        print('result of summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def generate_code(connectionId, requestId, chat, text, context, mode):
    if mode == 'py':    
        system = (
            """다음의 <context> tag안에는 질문과 관련된 python code가 있습니다. 주어진 예제를 참조하여 질문과 관련된 python 코드를 생성합니다. Assistant의 이름은 서연입니다. 결과는 <result> tag를 붙여주세요.
            
            <context>
            {context}
            </context>"""
        )
    elif mode == 'js':
        system = (
            """다음의 <context> tag안에는 질문과 관련된 node.js code가 있습니다. 주어진 예제를 참조하여 질문과 관련된 node.js 코드를 생성합니다. Assistant의 이름은 서연입니다. 결과는 <result> tag를 붙여주세요.
            
            <context>
            {context}
            </context>"""
        )
    
    human = "<context>{text}</context>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        isTyping(connectionId, requestId, "")  
        stream = chain.invoke(
            {
                "context": context,
                "text": text
            }
        )
        
        geenerated_code = readStreamMsg(connectionId, requestId, stream.content)
                              
        geenerated_code = stream.content        
        print('result of code generation: ', geenerated_code)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return geenerated_code

def summary_of_code(chat, code, mode):
    if mode == 'py':
        system = (
            "다음의 <article> tag에는 python code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    elif mode == 'js':
        system = (
            "다음의 <article> tag에는 node.js code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    else:
        system = (
            "다음의 <article> tag에는 code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    
    human = "<article>{code}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "code": code
            }
        )
        
        summary = result.content
        print('result of code summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def revise_question(connectionId, requestId, chat, query):    
    global history_length, token_counter_history    
    history_length = token_counter_history = 0
    
    isTyping(connectionId, requestId, "revising...")
        
    if isKorean(query)==True :      
        system = (
            ""
        )  
        human = """이전 대화를 참조하여, 다음의 <question>의 뜻을 명확히 하는 새로운 질문을 한국어로 생성하세요. 새로운 질문은 원래 질문의 중요한 단어를 반드시 포함합니다. 결과는 <result> tag를 붙여주세요.
        
        <question>            
        {question}
        </question>"""
        
    else: 
        system = (
            ""
        )
        human = """Rephrase the follow up <question> to be a standalone question. Put it in <result> tags.
        <question>            
        {question}
        </question>"""
            
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    # print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    # print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "history": history,
                "question": query,
            }
        )
        generated_question = result.content
        
        revised_question = generated_question[generated_question.find('<result>')+8:len(generated_question)-9] # remove <result> tag                   
        print('revised_question: ', revised_question)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':  
        chat_history = ""
        for dialogue_turn in history:
            #print('type: ', dialogue_turn.type)
            #print('content: ', dialogue_turn.content)
            
            dialog = f"{dialogue_turn.type}: {dialogue_turn.content}\n"            
            chat_history = chat_history + dialog
                
        history_length = len(chat_history)
        print('chat_history length: ', history_length)
        
        token_counter_history = 0
        if chat_history:
            token_counter_history = chat.get_num_tokens(chat_history)
            print('token_size of history: ', token_counter_history)
            
        sendDebugMessage(connectionId, requestId, f"새로운 질문: {revised_question}\n * 대화이력({str(history_length)}자, {token_counter_history} Tokens)을 활용하였습니다.")
            
    return revised_question    
    # return revised_question.replace("\n"," ")

def query_using_RAG_context(connectionId, requestId, chat, context, revised_question):    
    isTyping(connectionId, requestId, "generating...")  
    
    if isKorean(revised_question)==True:
        system = (
            """다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            
            <context>
            {context}
            </context>"""
        )
    else: 
        system = (
            """Here is pieces of context, contained in <context> tags. Provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer.
            
            <context>
            {context}
            </context>"""
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
                   
    chain = prompt | chat
    
    try: 
        isTyping(connectionId, requestId, "generating...") 
        stream = chain.invoke(
            {
                "context": context,
                "input": revised_question,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
        print('msg: ', msg)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    return msg
    
def use_multimodal(chat, img_base64, query):    
    if query == "":
        query = "그림에 대해 상세히 설명해줘."
    
    messages = [
        SystemMessage(content="답변은 500자 이내의 한국어로 설명해주세요."),
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    try: 
        result = chat.invoke(messages)
        
        summary = result.content
        print('result of code summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def extract_text(chat, img_base64):    
    query = "텍스트를 추출해서 utf8로 변환하세요. <result> tag를 붙여주세요."
    
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    try: 
        result = chat.invoke(messages)
        
        extracted_text = result.content
        print('result of text extraction from an image: ', extracted_text)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return extracted_text

# load documents from s3 
def load_document(file_type, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)
    
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
        
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text())
        contents = '\n'.join(texts)
        
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        from pptx import Presentation
        prs = Presentation(BytesIO(Byte_contents))

        texts = []
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = text + shape.text
            texts.append(text)
        contents = '\n'.join(texts)
        
    elif file_type == 'txt' or file_type == 'md':        
        contents = doc.get()['Body'].read().decode('utf-8')

    elif file_type == 'docx':
        Byte_contents = doc.get()['Body'].read()
            
        import docx
        doc_contents =docx.Document(BytesIO(Byte_contents))

        texts = []
        for i, para in enumerate(doc_contents.paragraphs):
            if(para.text):
                texts.append(para.text)
                # print(f"{i}: {para.text}")        
        contents = '\n'.join(texts)
            
    # print('contents: ', contents)
    new_contents = str(contents).replace("\n"," ") 
    print('length: ', len(new_contents))

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " ", ""],
        length_function = len,
    ) 

    texts = text_splitter.split_text(new_contents) 
                
    return texts

# load csv documents from s3
def load_csv_document(path, doc_prefix, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)

    lines = doc.get()['Body'].read().decode('utf-8').split('\n')   # read csv per line
    print('lins: ', len(lines))
        
    columns = lines[0].split(',')  # get columns
    #columns = ["Category", "Information"]  
    #columns_to_metadata = ["type","Source"]
    print('columns: ', columns)
    
    docs = []
    n = 0
    for row in csv.DictReader(lines, delimiter=',',quotechar='"'):
        # print('row: ', row)
        #to_metadata = {col: row[col] for col in columns_to_metadata if col in row}
        values = {k: row[k] for k in columns if k in row}
        content = "\n".join(f"{k.strip()}: {v.strip()}" for k, v in values.items())
        doc = Document(
            page_content=content,
            metadata={
                'name': s3_file_name,
                'page': n+1,
                'uri': path+doc_prefix+parse.quote(s3_file_name)
            }
            #metadata=to_metadata
        )
        docs.append(doc)
        n = n+1
    print('docs[0]: ', docs[0])

    return docs
    
def load_chat_history(userId, allowTime):
    dynamodb_client = boto3.client('dynamodb')

    response = dynamodb_client.query(
        TableName=callLogTableName,
        KeyConditionExpression='user_id = :userId AND request_time > :allowTime',
        ExpressionAttributeValues={
            ':userId': {'S': userId},
            ':allowTime': {'S': allowTime}
        }
    )
    # print('query result: ', response['Items'])

    for item in response['Items']:
        text = item['body']['S']
        msg = item['msg']['S']
        type = item['type']['S']

        if type == 'text' and text and msg:
            memory_chain.chat_memory.add_user_message(text)
            if len(msg) > MSG_LENGTH:
                memory_chain.chat_memory.add_ai_message(msg[:MSG_LENGTH])                          
            else:
                memory_chain.chat_memory.add_ai_message(msg) 
                                
def getAllowTime():
    d = datetime.datetime.now() - datetime.timedelta(days = 2)
    timeStr = str(d)[0:19]
    print('allow time: ',timeStr)

    return timeStr

def isTyping(connectionId, requestId, msg):    
    if not msg:
        msg = "typing a message..."
    msg_proceeding = {
        'request_id': requestId,
        'msg': msg,
        'status': 'istyping'
    }
    #print('result: ', json.dumps(result))
    sendMessage(connectionId, msg_proceeding)

def readStreamMsg(connectionId, requestId, stream):
    msg = ""
    if stream:
        for event in stream:
            # print('event: ', event)
            msg = msg + event

            result = {
                'request_id': requestId,
                'msg': msg,
                'status': 'proceeding'
            }
            #print('result: ', json.dumps(result))
            sendMessage(connectionId, result)
    # print('msg: ', msg)
    return msg

def get_ps_embedding():
    global selected_ps_embedding
    profile = priorty_search_embedding[selected_ps_embedding]
    bedrock_region =  profile['bedrock_region']
    model_id = profile['model_id']
    print(f'selected_ps_embedding: {selected_ps_embedding}, bedrock_region: {bedrock_region}, model_id: {model_id}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region, 
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    
    bedrock_ps_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = model_id
    )  
    
    selected_ps_embedding = selected_ps_embedding + 1
    if selected_ps_embedding == len(priorty_search_embedding):
        selected_ps_embedding = 0
    
    return bedrock_ps_embedding

def update_state_message(msg:str, config):
    print(msg)
    # print('config: ', config)
        
    requestId = config.get("configurable", {}).get("requestId", "")
    connectionId = config.get("configurable", {}).get("connectionId", "")
        
    isTyping(connectionId, requestId, msg)

def traslation(chat, text, input_language, output_language):
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        # print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def priority_search(query, relevant_docs, minSimilarity):
    excerpts = []
    for i, doc in enumerate(relevant_docs):
        #print('doc: ', doc)

        content = doc.page_content
        # print('content: ', content)

        excerpts.append(
            Document(
                page_content=content,
                metadata={
                    'name': doc.metadata['name'],
                    'url': doc.metadata['url'],
                    'from': doc.metadata['from'],
                    'order':i,
                    'score':0
                }
            )
        )
    #print('excerpts: ', excerpts)

    docs = []
    if len(excerpts):
        embeddings = get_ps_embedding()
        vectorstore_confidence = FAISS.from_documents(
            excerpts,  # documents
            embeddings  # embeddings
        )            
        rel_documents = vectorstore_confidence.similarity_search_with_score(
            query=query,
            k=len(relevant_docs)
        )
        
        for i, document in enumerate(rel_documents):
            print(f'## Document(priority_search) query: {query}, {i+1}: {document}')

            order = document[0].metadata['order']
            name = document[0].metadata['name']
            
            score = document[1]
            print(f"query: {query}, {order}: {name}, {score}")

            relevant_docs[order].metadata['score'] = int(score)

            if score < minSimilarity:
                docs.append(relevant_docs[order])    
        # print('selected docs: ', docs)
        
        # double check using translated query
        chat = get_chat()
        if isKorean(query):
            translated_query = traslation(chat, query, "Korean", "English")
        else:
            translated_query = traslation(chat, query, "English", "Korean")         
        print('translated_query: ', translated_query)
        
        rel_documents = vectorstore_confidence.similarity_search_with_score(
            query=translated_query,
            k=len(relevant_docs)
        )
        
        for i, document in enumerate(rel_documents):
            print(f'## Document(priority_search) query: {translated_query}, {i+1}: {document}')

            order = document[0].metadata['order']
            name = document[0].metadata['name']
            
            score = document[1]
            print(f"query: {translated_query}, {order}: {name}, {score}")

            relevant_docs[order].metadata['score'] = int(score)

            if score < minSimilarity:
                docs.append(relevant_docs[order])
        
        # check duplication of docs
        contentList = []
        updated_docs = []
        print('length of docs:', len(docs))
        for doc in docs:            
            if doc.page_content in contentList:
                print('duplicated!')
                continue
            contentList.append(doc.page_content)
            updated_docs.append(doc)
        print('length of updated_docs:', len(updated_docs))
        docs = updated_docs

    return docs
    
# get auth
region = os.environ.get('AWS_REGION', 'us-west-2')
print('region: ', region)
service = "aoss"  

credentials = boto3.Session().get_credentials()
awsauth = AWSV4SignerAuth(credentials, region, service)

os_client = OpenSearch(
    hosts = [{
        'host': opensearch_url.replace("https://", ""), 
        'port': 443
    }],
    http_auth=awsauth,
    use_ssl = True,
    verify_certs = True,
    connection_class=RequestsHttpConnection,
)

def is_not_exist(index_name):    
    print('index_name: ', index_name)
    
    #session = boto3.Session()
    #aoss_client = session.client(
    #    service_name="opensearchserverless"
    #)
    #response = aoss_client.list_collections()
    #print('response: ', response)
    
    if os_client.indices.exists(index_name):
        print('use exist index: ', index_name)    
        return False
    else:
        print('no index: ', index_name)
        return True

knowledge_base_id = ""
data_source_id = ""
def initiate_knowledge_base():
    global knowledge_base_id, data_source_id
    #########################
    # opensearch index
    #########################
    if(is_not_exist(vectorIndexName)):
        print(f"creating opensearch index... {vectorIndexName}")        
        body={
            'settings':{
                "index.knn": True,
                "index.knn.algo_param.ef_search": 512,
                'analysis': {
                    'analyzer': {
                        'my_analyzer': {
                            'char_filter': ['html_strip'], 
                            'tokenizer': 'nori',
                            'filter': ['nori_number','lowercase','trim','my_nori_part_of_speech'],
                            'type': 'custom'
                        }
                    },
                    'tokenizer': {
                        'nori': {
                            'decompound_mode': 'mixed',
                            'discard_punctuation': 'true',
                            'type': 'nori_tokenizer'
                        }
                    },
                    "filter": {
                        "my_nori_part_of_speech": {
                            "type": "nori_part_of_speech",
                            "stoptags": [
                                    "E", "IC", "J", "MAG", "MAJ",
                                    "MM", "SP", "SSC", "SSO", "SC",
                                    "SE", "XPN", "XSA", "XSN", "XSV",
                                    "UNA", "NA", "VSV"
                            ]
                        }
                    }
                },
            },
            'mappings': {
                'properties': {
                    'vector_field': {
                        'type': 'knn_vector',
                        'dimension': 1024,
                        'method': {
                            "name": "hnsw",
                            "engine": "faiss",
                            "parameters": {
                                "ef_construction": 512,
                                "m": 16
                            }
                        }                  
                    },
                    "AMAZON_BEDROCK_METADATA": {"type": "text", "index": False},
                    "AMAZON_BEDROCK_TEXT_CHUNK": {"type": "text"},
                }
            }
        }

        try: # create index
            response = os_client.indices.create(
                vectorIndexName,
                body=body
            )
            print('opensearch index was created:', response)

            # delay 3seconds
            time.sleep(5)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)                
            #raise Exception ("Not able to create the index")
            
    #########################
    # knowledge base
    #########################
    print('knowledge_base_name: ', knowledge_base_name)
    print('collectionArn: ', collectionArn)
    print('vectorIndexName: ', vectorIndexName)
    print('embeddingModelArn: ', embeddingModelArn)
    print('knowledge_base_role: ', knowledge_base_role)
    try: 
        client = boto3.client('bedrock-agent')         
        response = client.list_knowledge_bases(
            maxResults=10
        )
        print('(list_knowledge_bases) response: ', response)
        
        if "knowledgeBaseSummaries" in response:
            summaries = response["knowledgeBaseSummaries"]
            for summary in summaries:
                if summary["name"] == knowledge_base_name:
                    knowledge_base_id = summary["knowledgeBaseId"]
                    print('knowledge_base_id: ', knowledge_base_id)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)
                    
    if not knowledge_base_id:
        print('creating knowledge base...')        
        for atempt in range(3):
            try:
                response = client.create_knowledge_base(
                    name=knowledge_base_name,
                    description="Knowledge base based on OpenSearch",
                    roleArn=knowledge_base_role,
                    knowledgeBaseConfiguration={
                        'type': 'VECTOR',
                        'vectorKnowledgeBaseConfiguration': {
                            'embeddingModelArn': embeddingModelArn,
                            'embeddingModelConfiguration': {
                                'bedrockEmbeddingModelConfiguration': {
                                    'dimensions': 1024
                                }
                            }
                        }
                    },
                    storageConfiguration={
                        'type': 'OPENSEARCH_SERVERLESS',
                        'opensearchServerlessConfiguration': {
                            'collectionArn': collectionArn,
                            'fieldMapping': {
                                'metadataField': 'AMAZON_BEDROCK_METADATA',
                                'textField': 'AMAZON_BEDROCK_TEXT_CHUNK',
                                'vectorField': 'vector_field'
                            },
                            'vectorIndexName': vectorIndexName
                        }
                    }                
                )   
                print('(create_knowledge_base) response: ', response)
            
                if 'knowledgeBaseId' in response['knowledgeBase']:
                    knowledge_base_id = response['knowledgeBase']['knowledgeBaseId']
                    break
                else:
                    knowledge_base_id = ""    
            except Exception:
                    err_msg = traceback.format_exc()
                    print('error message: ', err_msg)
                    time.sleep(5)
                    print(f"retrying... ({atempt})")
                    #raise Exception ("Not able to create the knowledge base")       
                
    print(f"knowledge_base_name: {knowledge_base_name}, knowledge_base_id: {knowledge_base_id}")    
    
    #########################
    # data source      
    #########################
    data_source_name = s3_bucket  
    try: 
        response = client.list_data_sources(
            knowledgeBaseId=knowledge_base_id,
            maxResults=10
        )        
        print('(list_data_sources) response: ', response)
        
        if 'dataSourceSummaries' in response:
            for data_source in response['dataSourceSummaries']:
                print('data_source: ', data_source)
                if data_source['name'] == data_source_name:
                    data_source_id = data_source['dataSourceId']
                    print('data_source_id: ', data_source_id)
                    break    
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)
        
    if not data_source_id:
        print('creating data source...')  
        try:
            response = client.create_data_source(
                dataDeletionPolicy='DELETE',
                dataSourceConfiguration={
                    's3Configuration': {
                        'bucketArn': s3_arn,
                        'inclusionPrefixes': [ 
                            s3_prefix+'/',
                        ]
                    },
                    'type': 'S3'
                },
                description = f"S3 data source: {s3_bucket}",
                knowledgeBaseId = knowledge_base_id,
                name = data_source_name,
                vectorIngestionConfiguration={
                    'chunkingConfiguration': {
                        'chunkingStrategy': 'HIERARCHICAL',
                        'hierarchicalChunkingConfiguration': {
                            'levelConfigurations': [
                                {
                                    'maxTokens': 1500
                                },
                                {
                                    'maxTokens': 300
                                }
                            ],
                            'overlapTokens': 60
                        }
                    },
                    'parsingConfiguration': {
                        'bedrockFoundationModelConfiguration': {
                            'modelArn': parsingModelArn
                        },
                        'parsingStrategy': 'BEDROCK_FOUNDATION_MODEL'
                    }
                }
            )
            print('(create_data_source) response: ', response)
            
            if 'dataSource' in response:
                if 'dataSourceId' in response['dataSource']:
                    data_source_id = response['dataSource']['dataSourceId']
                    print('data_source_id: ', data_source_id)
                    
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            #raise Exception ("Not able to create the data source")
    
    print(f"data_source_name: {data_source_name}, data_source_id: {data_source_id}")
            
initiate_knowledge_base()

class GradeDocuments(BaseModel):
    """Binary score for relevance check on retrieved documents."""

    binary_score: str = Field(description="Documents are relevant to the question, 'yes' or 'no'")

def get_retrieval_grader(chat):
    system = """You are a grader assessing relevance of a retrieved document to a user question. \n 
    If the document contains keyword(s) or semantic meaning related to the question, grade it as relevant. \n
    Give a binary score 'yes' or 'no' score to indicate whether the document is relevant to the question."""

    grade_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system),
            ("human", "Retrieved document: \n\n {document} \n\n User question: {question}"),
        ]
    )
    
    structured_llm_grader = chat.with_structured_output(GradeDocuments)
    retrieval_grader = grade_prompt | structured_llm_grader
    return retrieval_grader

def grade_document_based_on_relevance(conn, question, doc, models, selected):     
    chat = get_multi_region_chat(models, selected)
    retrieval_grader = get_retrieval_grader(chat)
    score = retrieval_grader.invoke({"question": question, "document": doc.page_content})
    # print(f"score: {score}")
    
    grade = score.binary_score    
    if grade == 'yes':
        print("---GRADE: DOCUMENT RELEVANT---")
        conn.send(doc)
    else:  # no
        print("---GRADE: DOCUMENT NOT RELEVANT---")
        conn.send(None)
    
    conn.close()
                                    
def grade_documents_using_parallel_processing(question, documents):
    global selected_chat
    
    filtered_docs = []    

    processes = []
    parent_connections = []
    
    for i, doc in enumerate(documents):
        #print(f"grading doc[{i}]: {doc.page_content}")        
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        process = Process(target=grade_document_based_on_relevance, args=(child_conn, question, doc, multi_region_models, selected_chat))
        processes.append(process)

        selected_chat = selected_chat + 1
        if selected_chat == length_of_models:
            selected_chat = 0
    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        relevant_doc = parent_conn.recv()

        if relevant_doc is not None:
            filtered_docs.append(relevant_doc)

    for process in processes:
        process.join()
    
    #print('filtered_docs: ', filtered_docs)
    return filtered_docs

def grade_documents(question, documents):
    print("###### grade_documents ######")
    
    print("start grading...")
    print("grade_state: ", grade_state)
    
    if grade_state == "LLM":
        filtered_docs = []
        if multi_region == 'enable':  # parallel processing        
            filtered_docs = grade_documents_using_parallel_processing(question, documents)

        else:
            # Score each doc    
            chat = get_chat()
            retrieval_grader = get_retrieval_grader(chat)
            for i, doc in enumerate(documents):
                # print('doc: ', doc)
                print_doc(i, doc)
                
                score = retrieval_grader.invoke({"question": question, "document": doc.page_content})
                # print("score: ", score)
                
                grade = score.binary_score
                # print("grade: ", grade)
                # Document relevant
                if grade.lower() == "yes":
                    print("---GRADE: DOCUMENT RELEVANT---")
                    filtered_docs.append(doc)
                # Document not relevant
                else:
                    print("---GRADE: DOCUMENT NOT RELEVANT---")
                    # We do not include the document in filtered_docs
                    # We set a flag to indicate that we want to run web search
                    continue
    
    elif grade_state == "PRIORITY_SEARCH" and len(documents):
        filtered_docs = priority_search(question, documents, minDocSimilarity)
    else:  # OTHERS
        filtered_docs = documents

    return filtered_docs

def print_doc(i, doc):
    if len(doc.page_content)>=100:
        text = doc.page_content[:100]
    else:
        text = doc.page_content
            
    print(f"{i}: {text}, metadata:{doc.metadata}")

def get_docs_from_knowledge_base(documents):        
    relevant_docs = []
    for doc in documents:
        content = ""
        if doc.page_content:
            content = doc.page_content
        
        score = doc.metadata["score"]
        
        link = ""
        if "s3Location" in doc.metadata["location"]:
            link = doc.metadata["location"]["s3Location"]["uri"] if doc.metadata["location"]["s3Location"]["uri"] is not None else ""
            
            # print('link:', link)    
            pos = link.find(f"/{doc_prefix}")
            name = link[pos+len(doc_prefix)+1:]
            encoded_name = parse.quote(name)
            # print('name:', name)
            link = f"{path}{doc_prefix}{encoded_name}"
            
        elif "webLocation" in doc.metadata["location"]:
            link = doc.metadata["location"]["webLocation"]["url"] if doc.metadata["location"]["webLocation"]["url"] is not None else ""
            name = "WEB"

        url = link
        # print('url:', url)
        
        relevant_docs.append(
            Document(
                page_content=content,
                metadata={
                    'name': name,
                    'score': score,
                    'url': url,
                    'from': 'RAG'
                },
            )
        )    
    return relevant_docs
                
def get_answer_using_knowledge_base(chat, text, connectionId, requestId):    
    global reference_docs
    
    msg = ""
    top_k = numberOfDocs
    relevant_docs = []
    if knowledge_base_id:    
        isTyping(connectionId, requestId, "retrieving...")
        
        retriever = AmazonKnowledgeBasesRetriever(
            knowledge_base_id=knowledge_base_id, 
            retrieval_config={"vectorSearchConfiguration": {
                "numberOfResults": top_k,
                "overrideSearchType": "HYBRID"   # SEMANTIC
            }},
        )
        
        docs = retriever.invoke(text)
        # print('relevant_docs: ', relevant_docs)
        print('--> relevant_docs for knowledge base')
        for i, doc in enumerate(docs):
            print_doc(i, doc)
            
        relevant_docs = get_docs_from_knowledge_base(docs)
    
    # grading        
    isTyping(connectionId, requestId, "grading...")
    
    filtered_docs = grade_documents(text, relevant_docs)
    
    filtered_docs = check_duplication(filtered_docs) # duplication checker
            
    relevant_context = ""
    for i, document in enumerate(filtered_docs):
        print(f"{i}: {document}")
        if document.page_content:
            content = document.page_content
            
        relevant_context = relevant_context + content + "\n\n"
        
    print('relevant_context: ', relevant_context)

    msg = query_using_RAG_context(connectionId, requestId, chat, relevant_context, text)
    
    if len(filtered_docs):
        reference_docs += filtered_docs
            
    return msg
    
def traslation(chat, text, input_language, output_language):
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        # print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

####################### LangGraph #######################
# Chat Agent Executor
#########################################################

@tool
def get_current_time(format: str=f"%Y-%m-%d %H:%M:%S")->str:
    """Returns the current date and time in the specified format"""
    print("###### get_current_time ######")
    # f"%Y-%m-%d %H:%M:%S"
    
    format = format.replace('\'','')
    timestr = datetime.datetime.now(timezone('Asia/Seoul')).strftime(format)
    # print('timestr:', timestr)
    
    return timestr

def get_lambda_client(region):
    return boto3.client(
        service_name='lambda',
        region_name=region
    )

@tool 
def get_book_list(keyword: str) -> str:
    """
    Search book list by keyword and then return book list
    keyword: search keyword
    return: book list
    """
    print("###### get_book_list ######")
    
    keyword = keyword.replace('\'','')

    answer = ""
    url = f"https://search.kyobobook.co.kr/search?keyword={keyword}&gbCode=TOT&target=total"
    response = requests.get(url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.text, "html.parser")
        prod_info = soup.find_all("a", attrs={"class": "prod_info"})
        
        if len(prod_info):
            answer = "추천 도서는 아래와 같습니다.\n"
            
        for prod in prod_info[:5]:
            title = prod.text.strip().replace("\n", "")       
            link = prod.get("href")
            answer = answer + f"{title}, URL: {link}\n\n"
    
    return answer

@tool
def get_weather_info(city: str) -> str:
    """
    retrieve weather information by city name and then return weather statement.
    city: the name of city to retrieve
    return: weather statement
    """    
    print("###### get_weather_info ######")
    
    city = city.replace('\n','')
    city = city.replace('\'','')
    city = city.replace('\"','')
                
    chat = get_chat()
    if isKorean(city):
        place = traslation(chat, city, "Korean", "English")
        print('city (translated): ', place)
    else:
        place = city
        city = traslation(chat, city, "English", "Korean")
        print('city (translated): ', city)
        
    print('place: ', place)
    
    weather_str: str = f"{city}에 대한 날씨 정보가 없습니다."
    if weather_api_key: 
        apiKey = weather_api_key
        lang = 'en' 
        units = 'metric' 
        api = f"https://api.openweathermap.org/data/2.5/weather?q={place}&APPID={apiKey}&lang={lang}&units={units}"
        # print('api: ', api)
                
        try:
            result = requests.get(api)
            result = json.loads(result.text)
            print('result: ', result)
        
            if 'weather' in result:
                overall = result['weather'][0]['main']
                current_temp = result['main']['temp']
                min_temp = result['main']['temp_min']
                max_temp = result['main']['temp_max']
                humidity = result['main']['humidity']
                wind_speed = result['wind']['speed']
                cloud = result['clouds']['all']
                
                weather_str = f"{city}의 현재 날씨의 특징은 {overall}이며, 현재 온도는 {current_temp}도 이고, 최저온도는 {min_temp}도, 최고 온도는 {max_temp}도 입니다. 현재 습도는 {humidity}% 이고, 바람은 초당 {wind_speed} 미터 입니다. 구름은 {cloud}% 입니다."
                #weather_str = f"Today, the overall of {city} is {overall}, current temperature is {current_temp} degree, min temperature is {min_temp} degree, highest temperature is {max_temp} degree. huminity is {humidity}%, wind status is {wind_speed} meter per second. the amount of cloud is {cloud}%."            
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)                    
            # raise Exception ("Not able to request to LLM")    
        
    print('weather_str: ', weather_str)                            
    return weather_str

@tool
def search_by_tavily(keyword: str) -> str:
    """
    Search any information by keyword and then return the result as a string.
    keyword: search keyword which is greater than the minimum of 4 characters
    return: the information of keyword
    """    
    print("###### search_by_tavily ######")
    
    global reference_docs, selected_tavily
    
    print('tavily keyword: ', keyword)
    
    docs = []
    if selected_tavily != -1:
        selected_tavily = selected_tavily + 1
        if selected_tavily == len(tavily_api_key):
            selected_tavily = 0

        try:
            tavily_client = TavilyClient(api_key=tavily_api_key[selected_tavily])
            response = tavily_client.search(keyword, max_results=3)
            print('tavily response: ', response)
            
            print(f"--> tavily search result: {keyword}")
            for i, r in enumerate(response["results"]):
                content = r.get("content")
                print(f"{i}: {content}")

                name = r.get("title")
                if name is None:
                    name = 'WWW'
                    
                url = ""
                if "url" in r:
                    url = r.get("url")
            
                docs.append(
                    Document(
                        page_content=content,
                        metadata={
                            'name': name,
                            'url': url,
                            'from': 'tavily'
                        },
                    )
                )   
            
        except Exception as e:
            print('Exception: ', e)
        
        filtered_docs = grade_documents(keyword, docs)
        
        # duplication checker
        filtered_docs = check_duplication(filtered_docs)
        
    relevant_context = ""
    for i, document in enumerate(filtered_docs):
        print(f"{i}: {document}")
        if document.page_content:
            content = document.page_content
            
        relevant_context = relevant_context + content + "\n\n"        
    print('relevant_context: ', relevant_context)
        
    reference_docs += filtered_docs
        
    return relevant_context

@tool    
def search_by_knowledge_base(keyword: str) -> str:
    """
    Search technical information by keyword and then return the result as a string.
    keyword: search keyword
    return: the technical information of keyword
    """    
    print("###### search_by_knowledge_base ######")
    
    global reference_docs
    
    print('keyword: ', keyword)
    keyword = keyword.replace('\'','')
    keyword = keyword.replace('|','')
    keyword = keyword.replace('\n','')
    print('modified keyword: ', keyword)
    
    top_k = numberOfDocs
    relevant_docs = []
    if knowledge_base_id:    
        retriever = AmazonKnowledgeBasesRetriever(
            knowledge_base_id=knowledge_base_id, 
            retrieval_config={"vectorSearchConfiguration": {
                "numberOfResults": top_k,
                "overrideSearchType": "HYBRID"   # SEMANTIC
            }},
        )
        
        docs = retriever.invoke(keyword)
        # print('docs: ', docs)
        print('--> relevant_docs from knowledge base')
        for i, doc in enumerate(docs):
            print_doc(i, doc)
        
        relevant_docs = get_docs_from_knowledge_base(docs)

    # grading
    filtered_docs = grade_documents(keyword, relevant_docs)
    
    filtered_docs = check_duplication(filtered_docs) # duplication checker
            
    relevant_context = ""
    for i, document in enumerate(filtered_docs):
        print(f"{i}: {document}")
        if document.page_content:
            content = document.page_content
            
        relevant_context = relevant_context + content + "\n\n"        
    print('relevant_context: ', relevant_context)
    
    if len(filtered_docs):
        reference_docs += filtered_docs
        
    # print('reference_docs: ', reference_docs)
        
    return relevant_context

####################### LangGraph #######################
# Tool use
#########################################################
def buildToolAgent():
    class ToolState(TypedDict):
        # messages: Annotated[Sequence[BaseMessage], operator.add]
        messages: Annotated[list, add_messages]

    # tools = [get_current_time, get_book_list, get_weather_info, search_by_tavily, search_by_knowledge_base]    
    tools = [get_current_time, get_book_list, get_weather_info, search_by_tavily]    
    tool_node = ToolNode(tools)

    def call_model(state: ToolState, config):
        print("###### call_model ######")
        # print('state: ', state["messages"])
        
        update_state_message("thinking...", config)
        
        if isKorean(state["messages"][0].content)==True:
            system = (
                "당신의 이름은 서연이고, 질문에 친근한 방식으로 대답하도록 설계된 대화형 AI입니다."
                "상황에 맞는 구체적인 세부 정보를 충분히 제공합니다."
                "모르는 질문을 받으면 솔직히 모른다고 말합니다."
            )
        else: 
            system = (            
                "You are a conversational AI designed to answer in a friendly way to a question."
                "If you don't know the answer, just say that you don't know, don't try to make up an answer."
                "You will be acting as a thoughtful advisor."    
            )
            
        prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system),
                MessagesPlaceholder(variable_name="messages"),
            ]
        )
        
        chat = get_chat()             
        model = chat.bind_tools(tools)

        chain = prompt | model
            
        response = chain.invoke(state["messages"])
        print('call_model response: ', response.tool_calls)
        
        # state message
        if response.tool_calls:
            toolinfo = response.tool_calls[-1]            
            if toolinfo['type'] == 'tool_call':
                print('tool name: ', toolinfo['name'])                    
                update_state_message(f"calling... {toolinfo['name']}", config)
        
        return {"messages": [response]}
    
    def should_continue(state: ToolState) -> Literal["continue", "end"]:
        print("###### should_continue ######")
        messages = state["messages"]    
        # print('(should_continue) messages: ', messages)
            
        last_message = messages[-1]                
        if not last_message.tool_calls:
            next = "end"
        else:           
            next = "continue"             
        print(f"should_continue response: {next}")
        
        return next

    # build workflow
    workflow = StateGraph(ToolState)
    workflow.add_node("agent", call_model)
    workflow.add_node("action", tool_node)
    workflow.add_edge(START, "agent")
    workflow.add_conditional_edges(
        "agent",
        should_continue,
        {
            "continue": "action",
            "end": END,
        },
    )
    workflow.add_edge("action", "agent")

    return workflow.compile()

tool_app = buildToolAgent()    

def run_agent_executor(connectionId, requestId, query):        
    isTyping(connectionId, requestId, "thinking...")
    
    inputs = [HumanMessage(content=query)]
    config = {
        "recursion_limit": 50,
        "requestId": requestId,
        "connectionId": connectionId
    }
    
    message = ""
    
    for event in tool_app.stream({"messages": inputs}, config, stream_mode="values"):   
        # print('event: ', event)        
        message = event["messages"][-1]
        # print('message: ', message)

    msg = readStreamMsg(connectionId, requestId, message.content)
    
    # print('reference_docs: ', reference_docs)
    return msg

####################### LangGraph #######################
# Plan and Execute
#########################################################
def run_plan_and_exeucute(connectionId, requestId, query):
    class State(TypedDict):
        input: str
        plan: list[str]
        past_steps: Annotated[List[Tuple], operator.add]
        info: Annotated[List[Tuple], operator.add]
        answer: str

    class Plan(BaseModel):
        """List of steps as a json format"""

        steps: List[str] = Field(
            description="different steps to follow, should be in sorted order"
        )

    def get_planner():
        system = (
            "For the given objective, come up with a simple step by step plan."
            "This plan should involve individual tasks, that if executed correctly will yield the correct answer. Do not add any superfluous steps."
            "The result of the final step should be the final answer. Make sure that each step has all the information needed - do not skip steps."
        )
            
        planner_prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system),
                ("placeholder", "{messages}"),
            ]
        )
        
        chat = get_chat()   
        
        planner = planner_prompt | chat
        return planner
    
    def plan_node(state: State, config):
        print("###### plan ######")
        print('input: ', state["input"])
        
        update_state_message("planning...", config)
        
        inputs = [HumanMessage(content=state["input"])]

        planner = get_planner()
        response = planner.invoke({"messages": inputs})
        # print('response.content: ', response.content)
        
        for attempt in range(5):
            chat = get_chat()
            structured_llm = chat.with_structured_output(Plan, include_raw=True)
            info = structured_llm.invoke(response.content)
            print(f'attempt: {attempt}, info: {info}')
            
            if not info['parsed'] == None:
                parsed_info = info['parsed']
                # print('parsed_info: ', parsed_info)        
                print('steps: ', parsed_info.steps)                
                return {
                    "input": state["input"],
                    "plan": parsed_info.steps
                }
        
        print('parsing_error: ', info['parsing_error'])
        return {"plan": []}          

    def execute_node(state: State, config):
        print("###### execute ######")
        print('input: ', state["input"])
        plan = state["plan"]
        print('plan: ', plan) 
        
        update_state_message("executing...", config)
        
        plan_str = "\n".join(f"{i+1}. {step}" for i, step in enumerate(plan))
        print("plan_str: ", plan_str)
        
        task = plan[0]
        print('task: ', task)
                
        inputs = [HumanMessage(content=task)]
        output = tool_app.invoke({"messages": inputs}, config)
        print('executor output: ', output)
        
        transaction = []
        if "messages" in output:
            human = output["messages"][0]
            print('human: ', human.content)

            ai = output["messages"][-1]
            print('ai: ', ai.content)
            
            transaction = [HumanMessage(content=human.content), AIMessage(content=ai.content)]
            # print('transaction: ', transaction)
           
        return {
            "input": state["input"],
            "plan": state["plan"],
            "info": transaction,
            "past_steps": [task],
        }

    class Response(BaseModel):
        """Response to user."""
        response: str
        
    class Act(BaseModel):
        """Action to perform as a json format"""
        action: Union[Response, Plan] = Field(
            description="Action to perform. If you want to respond to user, use Response. "
            "If you need to further use tools to get the answer, use Plan."
        )

    def replan_node(state: State, config):
        print('#### replan ####')
        # print('state of replan node: ', state)
        
        update_state_message("replanning...", config)
        
        replanner_prompt = ChatPromptTemplate.from_template(
            "For the given objective, come up with a simple step by step plan."
            "This plan should involve individual tasks, that if executed correctly will yield the correct answer."
            "Do not add any superfluous steps."
            "The result of the final step should be the final answer."
            "Make sure that each step has all the information needed - do not skip steps."

            "Your objective was this:"
            "{input}"

            "Your original plan was this:"
            "{plan}"

            "You have currently done the follow steps:"
            "{past_steps}"

            "Update your plan accordingly."
            "If no more steps are needed and you can return to the user, then respond with that."
            "Otherwise, fill out the plan."
            "Only add steps to the plan that still NEED to be done. Do not return previously done steps as part of the plan."
        )
        
        chat = get_chat()
        replanner = replanner_prompt | chat
        
        output = replanner.invoke(state)
        print('replanner output: ', output.content)
        
        result = None
        for attempt in range(5):
            chat = get_chat()
            structured_llm = chat.with_structured_output(Act, include_raw=True)    
            info = structured_llm.invoke(output.content)
            print(f'attempt: {attempt}, info: {info}')
            
            if not info['parsed'] == None:
                result = info['parsed']
                print('replan result: ', result)
                break
                    
        if result == None:
            return {"response": "답을 찾지 못하였습니다. 다시 시도해주세요."}
        else:
            if isinstance(result.action, Response):  # "parsed":"Act(action=Response(response="
                return {
                    "response": result.action.response,
                    "info": [result.action.response]
                }
            else:  # "parsed":"Act(action=Plan(steps=
                return {"plan": result.action.steps}
        
    def should_end(state: State) -> Literal["continue", "end"]:
        print('#### should_end ####')
        # print('state: ', state)
        
        if "response" in state and state["response"]:
            print('response: ', state["response"])            
            next = "end"
        else:
            print('plan: ', state["plan"])
            next = "continue"
        print(f"should_end response: {next}")
        
        return next
        
    def final_answer(state: State) -> str:
        print('#### final_answer ####')
        
        # get final answer
        context = state['info']
        print('context: ', context)
        
        query = state['input']
        print('query: ', query)
        
        if isKorean(query)==True:
            system = (
                "Assistant의 이름은 서연이고, 질문에 대해 친절하게 답변하는 도우미입니다."
                "다음의 <context> tag안의 참고자료를 이용하여 질문에 대한 답변합니다."
                "답변의 이유를 풀어서 명확하게 설명합니다."
                "결과는 <result> tag를 붙여주세요."
                
                "<context>"
                "{context}"
                "</context>"
            )
        else: 
            system = (
                "Here is pieces of context, contained in <context> tags."
                "Provide a concise answer to the question at the end."
                "Explains clearly the reason for the answer."
                "If you don't know the answer, just say that you don't know, don't try to make up an answer."
                "Put it in <result> tags."
                
                "<context>"
                "{context}"
                "</context>"
            )
    
        human = "{input}"
        
        prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
        # print('prompt: ', prompt)
                    
        chat = get_chat()
        chain = prompt | chat
        
        try: 
            response = chain.invoke(
                {
                    "context": context,
                    "input": query,
                }
            )
            result = response.content
            output = result[result.find('<result>')+8:result.find('</result>')] # remove <result> tag
            print('output: ', output)
            
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)      
            
        return {"answer": output}  

    def buildPlanAndExecute():
        workflow = StateGraph(State)
        workflow.add_node("planner", plan_node)
        workflow.add_node("executor", execute_node)
        workflow.add_node("replaner", replan_node)
        workflow.add_node("final_answer", final_answer)
        
        workflow.set_entry_point("planner")
        workflow.add_edge("planner", "executor")
        workflow.add_edge("executor", "replaner")
        workflow.add_conditional_edges(
            "replaner",
            should_end,
            {
                "continue": "executor",
                "end": "final_answer",
            },
        )
        workflow.add_edge("final_answer", END)

        return workflow.compile()

    app = buildPlanAndExecute()    
    
    isTyping(connectionId, requestId, "")
    
    inputs = {"input": query}
    config = {
        "recursion_limit": 50,
        "requestId": requestId,
        "connectionId": connectionId
    }
    
    for output in app.stream(inputs, config):   
        for key, value in output.items():
            print(f"Finished: {key}")
            #print("value: ", value)            
    print('value: ', value)
    
    readStreamMsg(connectionId, requestId, value["answer"])
        
    return value["answer"]

####################### LangGraph #######################
# Plan and Execute
#########################################################
def get_llm(select):
    if multi_region == 'enable':
        if select>=len(multi_region_models):
            select -= len(multi_region_models)
        profile = multi_region_models[select]
    else:
        if select>=len(LLM_for_chat):
            select -= len(LLM_for_chat)
        profile = LLM_for_chat[select]
    
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    maxOutputTokens = 4096
    print(f'LLM: {select}, bedrock_region: {bedrock_region}, modelId: {modelId}')
                        
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }

    chat = ChatBedrock(   # new chat model
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )    
    
    select = select + 1
    if select == length_of_models:
        select = 0
    
    return chat, select

def solve_CSAT_Korean(connectionId, requestId, paragraph, question, question_plus, choices, idx, nth):    
    class State(TypedDict):
        plan: list[str]
        past_steps: Annotated[List[Tuple], operator.add]
        info: Annotated[List[Tuple], operator.add]
        paragraph: str
        question: str
        question_plus: str
        choices: list[str]
        answer: str
        select: int

    def plan_node(state: State, config):
        print("###### plan ######")
        print('paragraph: ', state["paragraph"])
        print('question: ', state["question"])
        print('question_plus: ', state["question_plus"])
                
        list_choices = ""
        choices = state["choices"]
        for i, choice in enumerate(choices):
            list_choices += f"({i+1}) {choice}\n"
        print('list_choices: ', list_choices)    
        
        idx = config.get("configurable", {}).get("idx")
        nth = config.get("configurable", {}).get("nth")
        update_state_message(f"planning... ({idx}-{nth})", config)
        
        system = (
            "당신은 복잡한 문제를 해결하기 위해 step by step plan을 생성하는 AI agent입니다."                
            
            "문제를 충분히 이해하고, 문제 해결을 위한 계획을 다음 형식으로 4단계 이하의 계획을 세웁니다."                
            "각 단계는 반드시 한줄의 문장으로 AI agent가 수행할 내용을 명확히 나타냅니다."
            "1. [질문을 해결하기 위한 단계]"
            "2. [질문을 해결하기 위한 단계]"
            "..."                
        )
        
        human = (
            "<paragraph> tag의 주어진 문장을 참조하여 <question> tag의 질문에 대한 적절한 답변을 <choice> tag안에서 선택하가 위한 단계별 계획을 세우세요."
            "단계별 계획에 <result> tag를 붙여주세요."
            
            "주어진 문장:"
            "<paragraph>"
            "{paragraph}"
            "</paragraph>"

            "질문:"
            "<question>"
            "{question}"
                            
            "{question_plus}"                
            "</question>"

            "선택지:"
            "<choices>"
            "{list_choices}"
            "</choices>"
        )
                            
        planner_prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system),
                ("human", human),
            ]
        )
        chat, select = get_llm(state["select"])
        planner = planner_prompt | chat
        response = planner.invoke({
            "paragraph": paragraph,
            "question": question,
            "question_plus": question_plus,
            "list_choices": list_choices
        })
        print('response.content: ', response.content)
        result = response.content
        output = result[result.find('<result>')+8:result.find('</result>')]
        
        plan = output.strip().replace('\n\n', '\n')
        planning_steps = plan.split('\n')
        print('planning_steps: ', planning_steps)
        
        return {
            "plan": planning_steps,
            "select": select
        }

    def execute_node(state: State, config):
        print("###### execute ######")
        plan = state["plan"]
        # print('plan: ', plan) 
        
        list_choices = ""
        choices = state["choices"]
        for i, choice in enumerate(choices):
            list_choices += f"({i+1}) {choice}\n"
        # print('list_choices: ', list_choices)    
                
        idx = config.get("configurable", {}).get("idx")
        nth = config.get("configurable", {}).get("nth")
        update_state_message(f"executing... ({idx}-{nth})", config)
        
        task = plan[0]
        print('task: ', task)                        
        
        context = ""
        for info in state['info']:
            if isinstance(info, HumanMessage):
                context += info.content+"\n"
            else:
                context += info.content+"\n\n"
        # print('context: ', context)
                        
        system = (
            "당신은 국어 수능문제를 푸는 일타강사입니다."
        )
        human = (
            "당신의 목표는 <paragraph> tag의 주어진 문장으로 부터 <question> tag의 주어진 질문에 대한 적절한 답변을 <choice> tag의 선택지에서 찾는것입니다."
            "<previous_result> tag에 있는 이전 단계의 결과를 참조하여, <task> tag의 실행 단계를 수행하고 적절한 답변을 구합니다."
            "문제를 풀이할 때 모든 선택지마다 근거를 주어진 문장에서 찾아 설명하세요."
            "선택지의 주요 단어들의 의미를 주어진 문장과 비교해서 꼼꼼히 차이점을 찾습니다."
            "질문에 대한 답을 선택지 중에 한 개만 골라서 대답해야 합니다."
            "최종 결과의 번호에 <result> tag를 붙여주세요."
            "최종 결과의 신뢰도를 1-5 사이의 숫자로 나타냅니다. 신뢰되는 <confidence> tag를 붙입니다."  
                                
            "주어진 문장:"
            "<paragraph>"
            "{paragraph}"
            "</paragraph>"
                
            "주어진 질문:"
            "<question>"
            "{question}"
                
            "{question_plus}"
            "</question>"
            
            "선택지:"
            "<choices>"
            "{list_choices}"
            "</choices>"
            
            "이전 단계의 결과"
            "<previous_result>"
            "{info}"
            "</previous_result>"

            "실행 단계:"
            "<task>"
            "{task}"
            "</task>"
        )
        prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system),
                ("human", human),
            ]
        )
        chat, select = get_llm(state["select"])
        chain = prompt | chat                        
        response = chain.invoke({
            "paragraph": state["paragraph"],
            "question": state["question"],
            "question_plus": state["question_plus"],
            "list_choices": list_choices,
            "info": context,
            "task": task
        })
        print('response.content: ', response.content)
        
        result = response.content
        confidence = result[result.find('<confidence>')+12:result.find('</confidence>')]
        print('confidence: ', confidence)
        
        transaction = [HumanMessage(content=task), AIMessage(content=result)]
        # print('transaction: ', transaction)
        
        if confidence == "5":
            plan = []
            answer = result
        else:
            plan = state["plan"]
            answer = ""
        
        return {
            "plan": plan,
            "info": transaction,
            "past_steps": [task],
            "answer": answer,
            "select": select
        }

    def replan_node(state: State, config):
        print('#### replan ####')
        # print('state of replan node: ', state)        
        print('past_steps: ', state["past_steps"])
                
        list_choices = ""
        choices = state["choices"]
        for i, choice in enumerate(choices):
            list_choices += f"({i+1}) {choice}\n"
        # print('list_choices: ', list_choices)    
        
        idx = config.get("configurable", {}).get("idx")
        nth = config.get("configurable", {}).get("nth")
        update_state_message(f"replanning... ({idx}-{nth})", config)
        
        if len(state["plan"])==0:
            return {"plan": []}
        
        system = (
            "당신은 복잡한 문제를 해결하기 위해 step by step plan을 생성하는 AI agent입니다."
        )        
        human = (
            "당신의 목표는 <paragraph> tag의 주어진 문장으로 부터 <question> tag의 주어진 질문에 대한 적절한 답변을 <choice> tag안에서 선택지에서 찾는것입니다."
            
            "주어진 문장:"
            "<paragraph>"
            "{paragraph}"
            "</paragraph>"
            
            "주어진 질문:"
            "<question>"
            "{question}"
            
            "{question_plus}"
            "</question>"
            
            "선택지:"
            "<list_choices>"
            "{list_choices}"
            "</list_choices>"
            
            "당신의 원래 계획은 아래와 같습니다." 
            "<original_plan>"                
            "{plan}"
            "</original_plan>"

            "완료한 단계는 아래와 같습니다."
            "<past_steps>"
            "{past_steps}"
            "</past_steps>"
            
            "당신은 <original_plan> tag의 원래 계획을 상황에 맞게 수정하세요."
            "계획에 아직 해야 할 단계만 추가하세요. 이전에 완료한 단계는 계획에 포함하지 마세요."                
            "수정된 계획에는 <plan> tag를 붙여주세요."
            "만약 더 이상 계획을 세우지 않아도 <question> tag의 주어진 질문에 답변할 있다면, 최종 결과로 <question>에 대한 답변을 <result> tag를 붙여 전달합니다."
            
            "수정된 계획의 형식은 아래와 같습니다."
            "각 단계는 반드시 한줄의 문장으로 AI agent가 수행할 내용을 명확히 나타냅니다."
            "1. [질문을 해결하기 위한 단계]"
            "2. [질문을 해결하기 위한 단계]"
            "..."         
        )                    
        replanner_prompt = ChatPromptTemplate.from_messages(
            [
                ("system", system),
                ("human", human),
            ]
        )        
        chat, select = get_llm(state["select"])
        replanner = replanner_prompt | chat        
        response = replanner.invoke({
            "paragraph": state["paragraph"],
            "question_plus": state["question_plus"],
            "question": state["question"],
            "list_choices": list_choices,
            "plan": state["plan"],
            "past_steps": state["past_steps"]
        })
        print('response.content: ', response.content)
        result = response.content
        
        print('find: ', result.find('<plan>'))
        
        if result.find('<plan>') == -1:
            return {"plan":[], "select":select, "answer":result}
        else:
            output = result[result.find('<plan>')+6:result.find('</plan>')]
            print('plan output: ', output)
            
            plans = output.strip().replace('\n\n', '\n')
            planning_steps = plans.split('\n')
            print('planning_steps: ', planning_steps)
        
            return {"plan": planning_steps, "select":select}
                
    def should_end(state: State) -> Literal["continue", "end"]:
        print('#### should_end ####')
        # print('state: ', state)
        
        plan = state["plan"]
        print('plan: ', plan)
        if len(plan)<=1:
            next = "end"
        else:
            next = "continue"
        print(f"should_end response: {next}")
        
        return next
        
    def final_answer(state: State, config) -> str:
        print('#### final_answer ####')
        
        idx = config.get("configurable", {}).get("idx")
        nth = config.get("configurable", {}).get("nth")
        update_state_message(f"finalizing... ({idx}-{nth})", config)
        
        if len(state["plan"])==0:
            return {"answer": state["answer"]}
        
        # get final answer
        info = state['info']
        
        context = ""
        for info in state['info']:
            if isinstance(info, HumanMessage):
                context += info.content+"\n"
            else:
                context += info.content+"\n\n"
        print('context: ', context)
                                
        print('paragraph: ', state["paragraph"])
        print('question: ', state["question"])
        print('question_plus: ', state["question_plus"])
        
        list_choices = ""
        choices = state["choices"]
        for i, choice in enumerate(choices):
            list_choices += f"({i+1}) {choice}\n"
        print('list_choices: ', list_choices)
        
        system = (
            "당신은 국어 수능문제를 푸는 일타강사입니다."                
        )        
        human = (
            "<context> tag에 있는 검토 결과를 활용하여, <paragraph> tag의 주어진 문장으로 부터 <question> tag의 주어진 질문에 대한 적절한 답변을 <choice> tag안에서 선택하려고 합니다."
            "질문에 대한 답을 선택지 중에 한 개만 골라서 대답해야 합니다."
            "답변의 이유를 풀어서 명확하게 설명합니다."
            "최종 결과의 번호에 <result> tag를 붙여주세요."  
            "최종 결과의 신뢰도를 1-5 사이의 숫자로 나타냅니다. 신뢰되는 <confidence> tag를 붙입니다."  
            
            "이전 단계에서 검토한 결과:"
            "<context>"
            "{context}"
            "</context>"
                            
            "주어진 문장:"
            "<paragraph>"
            "{paragraph}"
            "</paragraph>"

            "주어진 질문:"
            "<question>"
            "{question}"
                            
            "{question_plus}"                
            "</question>"

            "선택지:"
            "<choices>"
            "{list_choices}"
            "</choices>"       
        )
        prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
        chat, select = get_llm(state["select"])
        chain = prompt | chat        
        response = chain.invoke(
            {
                "context": context,
                "paragraph": state["paragraph"],                    
                "question": state["question"],
                "question_plus": state["question_plus"],
                "list_choices": list_choices
            }
        )
        result = response.content

        return {"answer":result, "select":select}  

    def buildPlanAndExecute():
        workflow = StateGraph(State)
        workflow.add_node("planner", plan_node)
        workflow.add_node("executor", execute_node)
        workflow.add_node("replaner", replan_node)
        workflow.add_node("final_answer", final_answer)
        
        workflow.set_entry_point("planner")
        workflow.add_edge("planner", "executor")
        workflow.add_edge("executor", "replaner")
        workflow.add_conditional_edges(
            "replaner",
            should_end,
            {
                "continue": "executor",
                "end": "final_answer",
            },
        )
        workflow.add_edge("final_answer", END)

        return workflow.compile()

    # run graph
    selection = idx+nth+2*nth
    if selection>=length_of_models:
        selection-=length_of_models
    print('selection: ', selection)
    
    app = buildPlanAndExecute()    
    
    isTyping(connectionId, requestId, "")
            
    inputs = {
        "question": question,
        "question_plus": question_plus,
        "paragraph": paragraph,
        "choices": choices,
        "select": selection
    }
    config = {
        "idx": idx,
        "nth": nth,
        "recursion_limit": 50,
        "requestId": requestId,
        "connectionId": connectionId
    }
    
    for output in app.stream(inputs, config):   
        for key, value in output.items():
            print(f"Finished: {key}")
            #print("value: ", value)            
    print('value: ', value)
    
    readStreamMsg(connectionId, requestId, value["answer"])
        
    return value["answer"]

#########################################################
contentList = []
def check_duplication(docs):
    global contentList
    length_original = len(docs)
    
    updated_docs = []
    print('length of relevant_docs:', len(docs))
    for doc in docs:            
        # print('excerpt: ', doc['metadata']['excerpt'])
            if doc.page_content in contentList:
                print('duplicated!')
                continue
            contentList.append(doc.page_content)
            updated_docs.append(doc)            
    length_updateed_docs = len(updated_docs)     
    
    if length_original == length_updateed_docs:
        print('no duplication')
    
    return updated_docs

def get_references(docs):
    reference = "\n\nFrom\n"
    for i, doc in enumerate(docs):
        page = ""
        if "page" in doc.metadata:
            page = doc.metadata['page']
            #print('page: ', page)            
        url = ""
        if "url" in doc.metadata:
            url = doc.metadata['url']
            #print('url: ', url)                
        name = ""
        if "name" in doc.metadata:
            name = doc.metadata['name']
            #print('name: ', name)     
           
        sourceType = ""
        if "from" in doc.metadata:
            sourceType = doc.metadata['from']
        #print('sourceType: ', sourceType)        
        
        #if len(doc.page_content)>=1000:
        #    excerpt = ""+doc.page_content[:1000]
        #else:
        #    excerpt = ""+doc.page_content
        excerpt = ""+doc.page_content
        # print('excerpt: ', excerpt)
        
        # for some of unusual case 
        #excerpt = excerpt.replace('"', '')        
        #excerpt = ''.join(c for c in excerpt if c not in '"')
        excerpt = re.sub('"', '', excerpt)
        print('excerpt(quotation removed): ', excerpt)
        
        if page:                
            reference = reference + f"{i+1}. {page}page in <a href={url} target=_blank>{name}</a>, {sourceType}, <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
        else:
            reference = reference + f"{i+1}. <a href={url} target=_blank>{name}</a>, {sourceType}, <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
    return reference

def translate_text(chat, text):
    global time_for_inference
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
        
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    if isKorean(text)==False :
        input_language = "English"
        output_language = "Korean"
    else:
        input_language = "Korean"
        output_language = "English"
                        
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':          
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
    
    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def check_grammer(chat, text):
    global time_for_inference
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
        
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장의 오류를 찾아서 설명하고, 오류가 수정된 문장을 답변 마지막에 추가하여 주세요."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Find the error in the sentence and explain it, and add the corrected sentence at the end of your answer."
        )
        
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    # print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        msg = result.content
        print('result of grammer correction: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':          
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
    
    return msg

def solve_problems_using_parallel_processing(connectionId, requestId, json_data):
    processes = []
    parent_connections = []
    
    total_idx = len(json_data)
    print('total_idx: ', total_idx)
    
    messages = []
    earn_score = 0
    for idx in range(total_idx):
        messages.append("")
        
    for idx, question_group in enumerate(json_data):
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
        
        print(f"idx:{idx} --> data:{question_group}")
        
        paragraph = question_group["paragraph"]
        print('paragraph: ', paragraph)
                    
        problems = question_group["problems"]
        print('problems: ', json.dumps(problems))
        
        process = Process(target=solve_problems, args=(child_conn, connectionId, requestId, paragraph, problems, idx, total_idx))
        processes.append(process)
        
    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        result = parent_conn.recv()
        print('result: ', result)
        
        idx = result["idx"]
        message = result["message"]
        score = result["score"]
        print(f"idx:{idx} --> socre: {score}, message:{message}")

        if message is not None:
            print('message: ', message)
            messages[idx] = message
            earn_score += score

    for process in processes:
        process.join()
            
    final_msg = ""   
    for message in messages:
        final_msg += message + '\n'
    
    print('earn_score: ', earn_score)
    print('final_msg: ', final_msg)
    
    return final_msg, earn_score

def solve_problems(conn, connectionId, requestId, paragraph, problems, idx, total_idx):
    message = f"{idx+1}/{total_idx}\n"
    
    earn_score = 0    
    for n, problem in enumerate(problems):
        print(f'--> problem[{n}]: {problem}')
    
        question = problem["question"]
        print('question: ', question)
        question_plus = ""
        if "question_plus" in problem:
            question_plus = problem["question_plus"]
            print('question_plus: ', question_plus)
        choices = problem["choices"]
        print('choices: ', choices)
        answer = problem["answer"]
        score = problem["score"]
            
        result = solve_CSAT_Korean(connectionId, requestId+str(idx)+str(n), paragraph, question, question_plus, choices, idx, n)
        print('result: ', result)
        
        output = result[result.find('<result>')+8:result.find('</result>')]
        print('output: ', output)
        
        if output.isnumeric():
            selected_answer = int(output)
            print('selected_answer: ', selected_answer)
        else:
            class Selection(BaseModel):
                select: int = Field(description="선택지의 번호")
            
            chat = get_chat()
            structured_llm = chat.with_structured_output(Selection, include_raw=True)
            
            info = structured_llm.invoke(output)
            selected_answer = 0
            for attempt in range(5):
                #print(f'attempt: {attempt}, info: {info}')
                if not info['parsed'] == None:
                    parsed_info = info['parsed']
                    #print('parsed_info: ', parsed_info)
                    selected_answer = parsed_info.select                    
                    print('selected_answer: ', selected_answer)
                    break

        if answer == selected_answer:
            message += f"{question} {selected_answer} (OK)\n"
            earn_score += int(score)
        else:
            message += f"{question} {selected_answer} (NOK, {answer}, -{score})\n"
            
    print('earn_score: ', earn_score) 
    print('message: ', message)
    
    sendResultMessage(connectionId, requestId+str(idx), message)
    
    conn.send({
        "idx": idx, 
        "message": message, 
        "score": earn_score
    })
    
    conn.close()

def solve_problems_in_paragraph(connectionId, requestId, paragraph, problems, idx, total_idx):
    message = f"{idx+1}/{total_idx}\n"
    
    earn_score = 0    
    for n, problem in enumerate(problems):
        print(f'--> problem[{n}]: {problem}')
    
        question = problem["question"]
        print('question: ', question)
        question_plus = ""
        if "question_plus" in problem:
            question_plus = problem["question_plus"]
            print('question_plus: ', question_plus)
        choices = problem["choices"]
        print('choices: ', choices)
        answer = problem["answer"]
        score = problem["score"]
            
        result = solve_CSAT_Korean(connectionId, requestId+str(idx)+str(n), paragraph, question, question_plus, choices, idx+n)
        print('result: ', result)
        
        output = result[result.find('<result>')+8:result.find('</result>')]
        print('output: ', output)
        
        if output.isnumeric():
            selected_answer = int(output)
            print('selected_answer: ', selected_answer)
        else:
            class Selection(BaseModel):
                select: int = Field(description="선택지의 번호")
            
            chat = get_chat()
            structured_llm = chat.with_structured_output(Selection, include_raw=True)
            
            info = structured_llm.invoke(output)
            selected_answer = 0
            for attempt in range(5):
                #print(f'attempt: {attempt}, info: {info}')
                if not info['parsed'] == None:
                    parsed_info = info['parsed']
                    #print('parsed_info: ', parsed_info)
                    selected_answer = parsed_info.select                    
                    print('selected_answer: ', selected_answer)
                    break

        if answer == selected_answer:
            message += f"{question} {selected_answer} (OK)\n"
            earn_score += int(score)
        else:
            message += f"{question} {selected_answer} (NOK, {answer}, -{score})\n"
            
    print('earn_score: ', earn_score)
    print('message: ', message)
    
    return {
        "idx:": idx, 
        "message": message, 
        "score": earn_score
    }

def getResponse(connectionId, jsonBody):
    userId  = jsonBody['user_id']
    # print('userId: ', userId)
    requestId  = jsonBody['request_id']
    # print('requestId: ', requestId)
    requestTime  = jsonBody['request_time']
    # print('requestTime: ', requestTime)
    type  = jsonBody['type']
    # print('type: ', type)
    body = jsonBody['body']
    # print('body: ', body)
    conv_type = jsonBody['conv_type']  # conversation type
    print('Conversation Type: ', conv_type)
    
    rag_type = ""
    if 'rag_type' in jsonBody:
        if jsonBody['rag_type']:
            rag_type = jsonBody['rag_type']  # RAG type
            print('rag_type: ', rag_type)
    
    global multi_region    
    if "multi_region" in jsonBody:
        multi_region = jsonBody['multi_region']
    print('multi_region: ', multi_region)
    
    global grade_state
    if "grade" in jsonBody:
        grade_state = jsonBody['grade']
    else:
        grade_state = 'LLM'
    print('grade_state: ', grade_state)
    
    global reference_docs, contentList
    reference_docs = []
    contentList = []
    
    global enableReference
    global map_chain, memory_chain, debugMessageMode
                 
    # Multi-LLM
    global selected_chat, length_of_models
    if multi_region == 'enable':
        length_of_models = len(multi_region_models)
        if selected_chat >= length_of_models:
            selected_chat = 0
        profile = multi_region_models[selected_chat]
        
    else:
        length_of_models = len(LLM_for_chat)
        if selected_chat >= length_of_models:
            selected_chat = 0    
        profile = LLM_for_chat[selected_chat]
        
    # print('length_of_models: ', length_of_models)    
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_chat: {selected_chat}, bedrock_region: {bedrock_region}, modelId: {modelId}')
      
    chat = get_chat()    
    bedrock_embedding = get_embedding()

    # allocate memory
    if userId in map_chain:  
        print('memory exist. reuse it!')        
        memory_chain = map_chain[userId]
        
    else: 
        print('memory does not exist. create new one!')
        memory_chain = ConversationBufferWindowMemory(memory_key="chat_history", output_key='answer', return_messages=True, k=10)
        map_chain[userId] = memory_chain
        
        allowTime = getAllowTime()
        load_chat_history(userId, allowTime)
        
    start = int(time.time())

    msg = ""
    reference = ""
    isControlMsg = False
    token_counter_input = 0
    time_for_revise = time_for_rag = time_for_priority_search = time_for_inference = 0
    
    if type == 'text' and body[:11] == 'list models':
        isControlMsg = True
        bedrock_client = boto3.client(
            service_name='bedrock',
            region_name=bedrock_region,
        )
        modelInfo = bedrock_client.list_foundation_models()    
        print('models: ', modelInfo)

        msg = f"The list of models: \n"
        lists = modelInfo['modelSummaries']
        
        for model in lists:
            msg += f"{model['modelId']}\n"
        
        msg += f"current model: {modelId}"
        print('model lists: ', msg)          

        sendResultMessage(connectionId, requestId, msg)  
    else:           
        text = body
        print('query: ', text)  
        querySize = len(text)
        textCount = len(text.split())
        print(f"query size: {querySize}, words: {textCount}")
        
        if type == 'text':
            if text == 'enableReference':
                enableReference = 'true'
                isControlMsg = True
                msg  = "Referece is enabled"
            elif text == 'disableReference':
                enableReference = 'false'
                isControlMsg = True
                msg  = "Reference is disabled"
            elif text == 'enableDebug':
                isControlMsg = True
                debugMessageMode = 'true'
                msg  = "Debug messages will be delivered to the client."
            elif text == 'disableDebug':
                isControlMsg = True
                debugMessageMode = 'false'
                msg = "Debug messages will not be delivered to the client."

            elif text == 'clearMemory':
                isControlMsg = True
                memory_chain.clear()
                map_chain[userId] = memory_chain
                    
                print('initiate the chat memory!')
                msg  = "The chat memory was intialized in this session."
            else:       
                if conv_type == 'normal':      # normal
                    msg = general_conversation(connectionId, requestId, chat, text)      
                    
                elif conv_type == 'qa-knowledge-base':   # RAG - Vector
                    print(f'rag_type: {rag_type}')
                    msg = get_answer_using_knowledge_base(chat, text, connectionId, requestId)
                
                elif conv_type == 'agent-executor':
                    msg = run_agent_executor(connectionId, requestId, text)
                
                elif conv_type == 'agent-executor-chat':
                    revised_question = revise_question(connectionId, requestId, chat, text)     
                    print('revised_question: ', revised_question)  
                    msg = run_agent_executor(connectionId, requestId, revised_question)
                
                elif conv_type == 'agent-plan-and-execute':  
                    msg = run_plan_and_exeucute(connectionId, requestId, text)        
                
                elif conv_type == 'agent-plan-and-execute-chat':
                    revised_question = revise_question(connectionId, requestId, chat, text)     
                    print('revised_question: ', revised_question)  
                    msg = run_plan_and_exeucute(connectionId, requestId, revised_question)        
                    
                elif conv_type == "translation":
                    msg = translate_text(chat, text) 
                elif conv_type == "grammar":
                    msg = check_grammer(chat, text)  
                                    
                # token counter
                if debugMessageMode=='true':
                    token_counter_input = chat.get_num_tokens(text)
                    token_counter_output = chat.get_num_tokens(msg)
                    print(f"token_counter: question: {token_counter_input}, answer: {token_counter_output}")
                    
                memory_chain.chat_memory.add_user_message(text)
                memory_chain.chat_memory.add_ai_message(msg)
                
                if reference_docs:
                    reference = get_references(reference_docs)
                        
        elif type == 'document':
            isTyping(connectionId, requestId, "")
            
            object = body
            file_type = object[object.rfind('.')+1:len(object)]            
            print('file_type: ', file_type)
            
            if file_type == 'csv':
                docs = load_csv_document(path, doc_prefix, object)
                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(chat, contexts)
                        
            elif file_type == 'pdf' or file_type == 'txt' or file_type == 'md' or file_type == 'pptx' or file_type == 'docx':
                texts = load_document(file_type, object)

                docs = []
                for i in range(len(texts)):
                    docs.append(
                        Document(
                            page_content=texts[i],
                            metadata={
                                'name': object,
                                # 'page':i+1,
                                'uri': path+doc_prefix+parse.quote(object)
                            }
                        )
                    )
                print('docs[0]: ', docs[0])    
                print('docs size: ', len(docs))

                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(chat, contexts)
                
            elif file_type == 'py' or file_type == 'js':
                s3r = boto3.resource("s3")
                doc = s3r.Object(s3_bucket, s3_prefix+'/'+object)
                
                contents = doc.get()['Body'].read().decode('utf-8')
                                
                msg = summary_of_code(chat, contents, file_type)          
            
            elif file_type == 'png' or file_type == 'jpeg' or file_type == 'jpg':
                print('multimodal: ', object)
                
                s3_client = boto3.client('s3') 
                    
                image_obj = s3_client.get_object(Bucket=s3_bucket, Key=s3_prefix+'/'+object)
                # print('image_obj: ', image_obj)
                
                image_content = image_obj['Body'].read()
                img = Image.open(BytesIO(image_content))
                
                width, height = img.size 
                print(f"width: {width}, height: {height}, size: {width*height}")
                
                isResized = False
                while(width*height > 5242880):                    
                    width = int(width/2)
                    height = int(height/2)
                    isResized = True
                    print(f"width: {width}, height: {height}, size: {width*height}")
                
                if isResized:
                    img = img.resize((width, height))
                
                buffer = BytesIO()
                img.save(buffer, format="PNG")
                img_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
                
                command  = jsonBody['command']
                print('command: ', command)
                
                # verify the image
                msg = use_multimodal(chat, img_base64, command)       
                
                # extract text from the image
                text = extract_text(chat, img_base64)
                extracted_text = text[text.find('<result>')+8:len(text)-9] # remove <result> tag
                print('extracted_text: ', extracted_text)
                if len(extracted_text)>10:
                    msg = msg + f"\n\n[추출된 Text]\n{extracted_text}\n"
                
                memory_chain.chat_memory.add_user_message(f"{object}에서 텍스트를 추출하세요.")
                memory_chain.chat_memory.add_ai_message(extracted_text)
            
            elif file_type == 'json':
                s3r = boto3.resource("s3")
                doc = s3r.Object(s3_bucket, s3_prefix+'/'+object)

                contents = doc.get()['Body'].read().decode('utf-8')   
                # print('contents: ', contents)
                
                json_data = json.loads(contents)
                    
                msg = ""  
                total_score = 0
                for question_group in json_data:
                    problems = question_group["problems"]
                    for problem in problems:
                        total_score += int(problem["score"])
                print('total_score: ', total_score)
                                        
                if multi_region_mode=="disable":
                    total_idx = len(json_data)+1
                    earn_score = 0
                    
                    for idx, question_group in enumerate(json_data[:2]):
                        paragraph = question_group["paragraph"]
                        print('paragraph: ', paragraph)
                        
                        problems = question_group["problems"]
                        print('problems: ', json.dumps(problems))
                        
                        result = solve_problems_in_paragraph(connectionId, requestId, paragraph, problems, idx, total_idx)
                        idx = result["idx"]
                        message = result["message"]
                        score = result["score"]
                        print('idx: ', idx)
                        print('message: ', message)
                        print('score: ', score)
                        
                        msg += message
                        earn_score += score
                        
                        msg += "\n\n"
                else:                    
                    msg, earn_score = solve_problems_using_parallel_processing(connectionId, requestId, json_data)
                    
                print('score: ', earn_score)
                msg += f"\n점수: {earn_score}점 / {total_score}점\n"
                
                sendResultMessage(connectionId, requestId, msg)  
                
            else:
                msg = "uploaded file: "+object
                
            # trigger sync of data source
            if knowledge_base_id and data_source_id:
                try:
                    client = boto3.client('bedrock-agent')
                    response = client.start_ingestion_job(
                        knowledgeBaseId=knowledge_base_id,
                        dataSourceId=data_source_id
                    )
                    print('(start_ingestion_job) response: ', response)
                except Exception:
                    err_msg = traceback.format_exc()
                    print('error message: ', err_msg)
                                                        
        sendResultMessage(connectionId, requestId, msg+reference)
        # print('msg+reference: ', msg+reference)

        elapsed_time = time.time() - start
        print("total run time(sec): ", elapsed_time)
                               
        item = {    # save dialog
            'user_id': {'S':userId},
            'request_id': {'S':requestId},
            'request_time': {'S':requestTime},
            'type': {'S':type},
            'body': {'S':body},
            'msg': {'S':msg+reference}
        }
        print('item: ', item)
        client = boto3.client('dynamodb')
        try:
            resp =  client.put_item(TableName=callLogTableName, Item=item)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            raise Exception ("Not able to write into dynamodb")        
        #print('resp, ', resp)

    if debugMessageMode=='true' and isControlMsg==False: 
        statusMsg = f"\n\n[통계]\nRegion: {bedrock_region}\nModelId: {modelId}\n"
        if token_counter_input:
            statusMsg = statusMsg + f"Question: {str(len(text))}자 / {token_counter_input}토큰\nAnswer: {str(len(msg))}자 / {token_counter_output}토큰\n"
            
        if history_length:
            statusMsg = statusMsg + f"History: {str(history_length)}자 / {token_counter_history}토큰\n"
            
        statusMsg = statusMsg + f"Time(초): "
        if time_for_revise != 0:
            statusMsg = statusMsg + f"{time_for_revise:.2f}(Revise), "
        if time_for_rag != 0:
            statusMsg = statusMsg + f"{time_for_rag:.2f}(RAG), "
        if time_for_priority_search != 0:
            statusMsg = statusMsg + f"{time_for_priority_search:.2f}(Priority) "
        if time_for_inference != 0:
            statusMsg = statusMsg + f"{time_for_inference:.2f}(Inference), "
        statusMsg = statusMsg + f"{elapsed_time:.2f}(전체)"
            
        sendResultMessage(connectionId, requestId, msg+reference+statusMsg)

    return msg, reference

def lambda_handler(event, context):
    # print('event: ', event)
    
    msg = ""
    if event['requestContext']: 
        connectionId = event['requestContext']['connectionId']        
        routeKey = event['requestContext']['routeKey']
        
        if routeKey == '$connect':
            print('connected!')
        elif routeKey == '$disconnect':
            print('disconnected!')
        else:
            body = event.get("body", "")
            #print("data[0:8]: ", body[0:8])
            if body[0:8] == "__ping__":
                # print("keep alive!")                
                sendMessage(connectionId, "__pong__")
            else:
                print('connectionId: ', connectionId)
                print('routeKey: ', routeKey)
        
                jsonBody = json.loads(body)
                print('request body: ', json.dumps(jsonBody))

                requestId  = jsonBody['request_id']
                try:
                    msg, reference = getResponse(connectionId, jsonBody)

                    print('msg+reference: ', msg+reference)
                                        
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)

                    sendErrorMessage(connectionId, requestId, err_msg)    
                    raise Exception ("Not able to send a message")

    return {
        'statusCode': 200
    }
